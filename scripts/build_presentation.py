"""Build the project presentation as a PPTX file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Color palette — Ocean/Satellite theme
NAVY = RGBColor(0x1E, 0x27, 0x61)
TEAL = RGBColor(0x06, 0x5A, 0x82)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK = RGBColor(0x21, 0x29, 0x5C)
GREEN = RGBColor(0x41, 0xAE, 0x76)
RED = RGBColor(0xEF, 0x65, 0x48)
ORANGE = RGBColor(0xF0, 0xB2, 0x7A)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)

TITLE_FONT = "Georgia"
BODY_FONT = "Calibri"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_text(slide, text, left, top, width, height, font_size=36, color=WHITE, bold=True, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = TITLE_FONT
    p.alignment = align
    return txBox


def add_body_text(slide, text, left, top, width, height, font_size=14, color=DARK, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = BODY_FONT
    return txBox


def add_bullets(slide, items, left, top, width, height, font_size=13, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = BODY_FONT
        p.space_after = Pt(4)
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        # Remove existing bullets first
        for existing in pPr.findall(qn('a:buChar')):
            pPr.remove(existing)
        for existing in pPr.findall(qn('a:buNone')):
            pPr.remove(existing)
        pPr.append(buChar)
    return txBox


def add_table(slide, headers, rows, left, top, width, height):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = BODY_FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = DARK
                paragraph.font.name = BODY_FONT
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return table_shape


def add_section_header(slide, number, title, left=0.5, top=0.15, width=9):
    add_body_text(slide, f"0{number}" if number < 10 else str(number), left, top, 0.5, 0.4, font_size=28, color=TEAL, bold=True)
    add_title_text(slide, title, left + 0.5, top + 0.05, width, 0.5, font_size=24, color=DARK, bold=True)
    # Separator line
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(left), Inches(top + 0.55), Inches(9), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # ===== SLIDE 1: Title =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, NAVY)
    add_title_text(slide, "Road Segmentation\nfrom Satellite Imagery", 0.8, 0.8, 8, 2, font_size=40, color=WHITE)
    add_body_text(slide, "Hudhud ML Engineer \u2014 Map Gen AI Team", 0.8, 2.8, 8, 0.5, font_size=18, color=ICE)
    add_body_text(slide, "Mina Essam", 0.8, 3.5, 8, 0.4, font_size=16, color=ICE)
    add_body_text(slide, "wandb.ai/minaessam/road-segmentation  \u2022  github.com/minaessam2015/road-segmentation", 0.8, 4.5, 8, 0.4, font_size=11, color=ICE)

    # ===== SLIDE 2: Problem Statement =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, 1, "Problem Statement")
    add_bullets(slide, [
        "Extract road pixels from satellite images (DeepGlobe dataset, 1024\u00d71024 RGB)",
        "Binary semantic segmentation: road vs background",
        "Severe class imbalance: roads are only ~4% of pixels",
        "Thin, elongated structures sensitive to small prediction errors",
        "Topology matters: a broken intersection makes the road network unusable",
        "Deliverables: trained model + HTTP inference API",
    ], 0.5, 1.0, 9, 3.5)

    # ===== SLIDE 3: Dataset & EDA =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, 2, "Dataset & EDA Findings")
    add_bullets(slide, [
        "6,226 training pairs (1024\u00d71024 RGB + binary road mask)",
        "Only train split has ground truth \u2014 valid/test have no masks",
        "Mean road coverage: 4.3%, median: 2.15% \u2014 severe class imbalance",
        "All masks binary, all dimensions consistent, no corrupted files",
        "Road widths: ~10px (rural) to ~60px (highway)",
    ], 0.5, 1.0, 5, 2.5)
    add_body_text(slide, "EDA-Driven Decisions", 6, 1.0, 3.5, 0.3, font_size=14, color=TEAL, bold=True)
    add_bullets(slide, [
        "Loss: BCE + Dice (not plain BCE)",
        "Metric: IoU primary (not pixel accuracy)",
        "Split: 85/15 stratified by coverage",
        "Input: 512\u21921024 progressive training",
        "Augmentations: flips, rotations, color jitter",
    ], 6, 1.4, 3.5, 2.5, font_size=12)

    # ===== SLIDE 4: Architecture Selection =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, 3, "Architecture Selection")
    add_body_text(slide, "Preliminary comparison (1000 samples, 15 epochs)", 0.5, 0.9, 9, 0.3, font_size=12, color=MID_GRAY)
    add_table(slide,
        ["Model", "IoU", "Dice", "Params", "Time"],
        [
            ["U-Net + ResNet34", "0.528", "0.691", "24.4M", "10 min"],
            ["LinkNet + ResNet34", "0.507", "0.673", "21.8M", "9 min"],
            ["SegFormer + MIT-B2", "0.485", "0.653", "24.7M", "15 min"],
            ["DeepLabV3+ + ResNet50", "0.484", "\u2014", "26.7M", "10 min"],
        ],
        0.5, 1.3, 9, 1.8
    )
    add_body_text(slide, "Final choice: UNet++ + EfficientNet-B4 at 1024\u00d71024", 0.5, 3.5, 9, 0.4, font_size=16, color=TEAL, bold=True)
    add_bullets(slide, [
        "UNet++ dense skip connections preserve thin road features",
        "EfficientNet-B4: NAS-optimized, better accuracy/param ratio than ResNet",
        "Full 1024\u00d71024 resolution \u2014 no downsampling, no detail loss",
    ], 0.5, 4.0, 9, 1.5, font_size=12)

    # ===== SLIDE 5: Training Pipeline =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, 4, "Training Pipeline")
    add_bullets(slide, [
        "YAML-configurable: 9 configs for different architectures",
        "Mixed precision (AMP), gradient accumulation, 2-phase encoder freeze/unfreeze",
        "BCE + Dice compound loss (handles 4% class imbalance)",
        "Boundary-weighted loss fine-tuning (5\u00d7 weight on road edges)",
        "Cosine annealing LR with warm-up, AdamW optimizer",
        "W&B experiment tracking: real-time metrics, prediction images, confidence histograms",
        "Colab-friendly: checkpoints on Google Drive, auto-resume on disconnect",
        "Stratified train/val split (85/15) with separate seed for reproducibility",
    ], 0.5, 1.0, 9, 4)

    # ===== SLIDE 6: Final Results =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_title_text(slide, "Final Model Results", 0.8, 0.3, 8, 0.6, font_size=28, color=WHITE)
    # Big number callouts
    add_title_text(slide, "0.7016", 1, 1.2, 3, 1, font_size=54, color=GREEN, align=PP_ALIGN.CENTER)
    add_body_text(slide, "Validation IoU", 1, 2.2, 3, 0.4, font_size=16, color=ICE)
    add_title_text(slide, "0.8246", 5.5, 1.2, 3, 1, font_size=54, color=GREEN, align=PP_ALIGN.CENTER)
    add_body_text(slide, "Validation Dice", 5.5, 2.2, 3, 0.4, font_size=16, color=ICE)
    add_bullets(slide, [
        "Architecture: UNet++ + EfficientNet-B4",
        "Input: 1024\u00d71024 (full resolution)",
        "Training: 50 epochs, 5292 train / 934 val samples",
        "Boundary loss fine-tuning: minimal additional gain (model well-calibrated)",
    ], 0.8, 3.0, 8.5, 2, font_size=13, color=ICE)

    # ===== SLIDE 7: Data Leakage =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, 5, "Data Integrity \u2014 Leakage Analysis")
    add_bullets(slide, [
        "DeepGlobe tiles cut from larger images \u2014 adjacent tiles may leak across splits",
        "Analyzed top-K IoU samples using ResNet50 + CLIP ViT-B/32 embeddings",
        "ResNet threshold 0.93: 39 leaky pairs flagged",
        "CLIP threshold 0.96: 104 leaky pairs flagged",
        "Combined (union): 123 unique leaky val samples (13%)",
    ], 0.5, 1.0, 9, 2.5)
    add_table(slide,
        ["Metric", "Original", "Cleaned", "Inflation"],
        [
            ["Mean IoU", "0.694", "0.675", "1.9%"],
        ],
        2, 3.5, 5, 0.7
    )
    add_body_text(slide, "Conclusion: metrics slightly optimistic \u2014 report both values for transparency", 0.5, 4.5, 9, 0.4, font_size=13, color=MID_GRAY)

    # ===== SLIDE 8: Error Analysis =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, 6, "Error Analysis & Calibration")
    add_bullets(slide, [
        "Per-sample IoU distribution: mean 0.694, median 0.71",
        "FP vs FN balance: roughly equal \u2014 model not biased toward over/under-prediction",
        "Performance by road coverage: sparse roads (<2%) are harder",
        "Performance by road width: thin roads (<6px) have lower IoU",
        "Spatial error heatmap: no edge artifacts detected",
        "Calibration: reliability diagram + ECE computed",
        "Temperature scaling tested for probability correction",
        "Findings: model is well-calibrated, post-processing should be minimal",
    ], 0.5, 1.0, 9, 4)

    # ===== SLIDE 9: Post-Processing =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, 7, "Post-Processing Ablation")
    add_table(slide,
        ["Step", "IoU", "Delta", "Verdict"],
        [
            ["Raw threshold (0.5)", "0.7016", "baseline", "BASELINE"],
            ["Optimal threshold (0.45)", "0.7017", "+0.01%", "Neutral"],
            ["+ Component filter", "0.7019", "+0.02%", "Neutral"],
            ["+ Morph. closing", "0.7012", "-0.04%", "Neutral"],
            ["+ Morph. opening", "0.7012", "+0.00%", "Neutral"],
            ["+ Gap bridging", "0.7012", "+0.00%", "Neutral"],
            ["TTA (6-fold)", "0.7014", "+0.65%*", "Include"],
        ],
        0.5, 1.0, 9, 2.8
    )
    add_body_text(slide, "* TTA gain measured on 100-sample subset", 0.5, 4.0, 9, 0.3, font_size=11, color=MID_GRAY)
    add_body_text(slide, "Model produces clean predictions \u2014 minimal post-processing needed", 0.5, 4.5, 9, 0.4, font_size=14, color=TEAL, bold=True)

    # ===== SLIDE 10: Optimization =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, 8, "Model Optimization")
    add_table(slide,
        ["Variant", "Latency", "Speedup", "IoU", "Size"],
        [
            ["PyTorch FP32", "145.6 ms", "1.00\u00d7", "0.6835", "\u2014"],
            ["PyTorch FP16/AMP", "95.5 ms", "1.52\u00d7", "0.6835", "\u2014"],
            ["ONNX FP32", "135.7 ms", "1.07\u00d7", "0.6835", "80 MB"],
            ["ONNX FP16", "78.3 ms", "1.86\u00d7", "0.6835", "40 MB"],
            ["ONNX INT8 (CPU)", "\u2014", "\u2014", "\u2014", "21 MB"],
        ],
        0.5, 1.0, 9, 2.2
    )
    add_title_text(slide, "Zero IoU loss across all variants", 0.5, 3.5, 9, 0.5, font_size=20, color=GREEN, bold=True)
    add_body_text(slide, "ONNX FP16 recommended for GPU production (1.86\u00d7 speedup, half the size)", 0.5, 4.1, 9, 0.4, font_size=14, color=DARK)

    # ===== SLIDE 11: API =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, 9, "Inference API")
    add_bullets(slide, [
        "FastAPI with 4 endpoints:",
        "    POST /api/v1/segment \u2014 upload image, get road mask + metadata",
        "    GET /health \u2014 model status, disk space, GPU memory",
        "    GET /api/v1/model-info \u2014 model version and config",
        "    GET /metrics \u2014 Prometheus-compatible metrics",
        "",
        "Supports PyTorch checkpoint and ONNX backends",
        "Input validation: format checks, size limits, helpful error messages",
        "Optional GeoJSON polyline output (skeletonization + graph extraction)",
        "Configurable threshold via query parameter",
    ], 0.5, 1.0, 5, 3.5)
    add_body_text(slide, "Observability", 6, 1.0, 3.5, 0.3, font_size=14, color=TEAL, bold=True)
    add_bullets(slide, [
        "Structured JSON logging",
        "X-Request-ID tracing",
        "Inference audit trail",
        "Prometheus /metrics endpoint",
        "Latency p50/p95/p99",
        "Error rate tracking",
    ], 6, 1.4, 3.5, 2.5, font_size=12)

    # ===== SLIDE 12: CI/CD =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, 10, "CI/CD & Deployment")
    add_body_text(slide, "GitHub Actions Pipeline", 0.5, 0.9, 5, 0.3, font_size=14, color=TEAL, bold=True)
    add_bullets(slide, [
        "Lint (ruff) \u2192 Unit tests (69) \u2192 Security scan",
        "\u2192 Docker build \u2192 E2E test with real model \u2192 Push to GHCR",
    ], 0.5, 1.3, 5, 1.0, font_size=12)
    add_body_text(slide, "Docker", 0.5, 2.3, 5, 0.3, font_size=14, color=TEAL, bold=True)
    add_bullets(slide, [
        "Multi-stage build (builder + slim runtime)",
        "Health check, volume mounts for models",
        "docker-compose for single-command startup",
    ], 0.5, 2.6, 5, 1.0, font_size=12)
    add_body_text(slide, "Model Registry", 6, 0.9, 3.5, 0.3, font_size=14, color=TEAL, bold=True)
    add_bullets(slide, [
        "W&B Artifacts: checkpoint + 3 ONNX variants",
        "python scripts/download_model.py",
        "Versioned, tagged, reproducible",
    ], 6, 1.3, 3.5, 1.0, font_size=12)
    add_body_text(slide, "69 Tests", 6, 2.3, 3.5, 0.3, font_size=14, color=TEAL, bold=True)
    add_bullets(slide, [
        "Config, transforms, models, losses",
        "Post-processing pipeline",
        "API endpoints (mocked + real model)",
        "End-to-end: PyTorch \u2192 ONNX \u2192 API",
    ], 6, 2.6, 3.5, 1.2, font_size=12)

    # ===== SLIDE 13: Production =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, 11, "Production Considerations")
    # 2x2 grid
    topics = [
        ("Scalability", [
            "Stateless service \u2192 horizontal scaling",
            "Task queues for batch processing",
            "Tiling with overlap for large images",
        ]),
        ("Latency", [
            "ONNX FP16: 78ms on T4 GPU",
            "TensorRT potential: additional 2-3\u00d7",
            "INT8 for CPU: 21 MB model",
        ]),
        ("Model Lifecycle", [
            "W&B registry versioning",
            "A/B testing, canary deploys",
            "One-command rollback",
        ]),
        ("Data Drift", [
            "Monitor input RGB statistics",
            "Track prediction distributions",
            "Calibration ECE over time",
        ]),
    ]
    positions = [(0.5, 1.0), (5.2, 1.0), (0.5, 3.0), (5.2, 3.0)]
    for (title, items), (x, y) in zip(topics, positions):
        add_body_text(slide, title, x, y, 4, 0.3, font_size=14, color=TEAL, bold=True)
        add_bullets(slide, items, x, y + 0.35, 4.2, 1.5, font_size=11)

    # ===== SLIDE 14: Notebooks =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_header(slide, 12, "Analysis Pipeline \u2014 7 Notebooks")
    add_table(slide,
        ["#", "Notebook", "Purpose"],
        [
            ["01", "EDA", "Data exploration, class imbalance, EDA-driven decisions"],
            ["02", "Model Comparison", "4-architecture quick benchmark (1000 samples)"],
            ["03", "Training", "Full training with W&B tracking + Drive checkpoints"],
            ["04", "Data Leakage", "Train/val similarity (ResNet + CLIP embeddings)"],
            ["05", "Error Analysis", "Failure patterns, calibration, spatial errors"],
            ["06", "Post-Processing", "Per-step ablation study (threshold, morphology, TTA)"],
            ["07", "Optimization", "ONNX export + speed/accuracy benchmark"],
        ],
        0.5, 1.0, 9, 3.0
    )

    # ===== SLIDE 15: Summary =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_title_text(slide, "Summary", 0.8, 0.3, 8, 0.6, font_size=32, color=WHITE)
    add_bullets(slide, [
        "End-to-end system: EDA \u2192 training \u2192 analysis \u2192 optimization \u2192 API \u2192 deployment",
        "IoU 0.7016 on DeepGlobe (competitive with published results)",
        "Production-ready: ONNX FP16 (78ms, zero accuracy loss), Docker, CI/CD",
        "Thorough analysis: data leakage, calibration, error patterns, post-processing ablation",
        "69 automated tests including end-to-end with real model weights",
        "All experiments tracked on W&B, models versioned in artifact registry",
    ], 0.8, 1.2, 8.5, 3, font_size=14, color=ICE)
    add_body_text(slide, "github.com/minaessam2015/road-segmentation", 0.8, 4.5, 8, 0.4, font_size=13, color=ICE)

    # Save
    output_path = "docs/presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
